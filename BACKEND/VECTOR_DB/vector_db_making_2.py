import os
import json
import certifi
import pymupdf
from pptx import Presentation
from dotenv import load_dotenv
from typing import TypedDict, Annotated
from langchain_openai import ChatOpenAI
from langchain_core.prompts import PromptTemplate

load_dotenv()
os.environ["SSL_CERT_FILE"] = certifi.where()


# -------------------------------------------------------------------
# PRIMARY FUNCTION TO CONVERT PDFS INTO VECTOR DB SCHEMA ORIENTED JSON
# --------------------------------------------------------------------


# ------------- PDF LOADER ----------------------

def extract_pdf_pages(pdf_path: str) -> list[str]:
    """
    Extracts clean text from each page of a PDF.
    Returns a list of page texts.
    """
    doc = pymupdf.open(pdf_path)
    pages = []

    for page in doc:
        text = page.get_text("text")
        
        if text.strip():
            pages.append(text)

    return pages


# ------------- PPT LOADER ----------------------

def extract_ppt_text(ppt_path: str) -> str:
    prs = Presentation(ppt_path)
    slides_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slides_text.append(shape.text)

    return "\n".join(slides_text)


# ---------------- TEXT LOADER ------------------------------
def extract_txt_text(txt_path: str, encoding: str = "utf-8") -> str:
    """
    Extracts and returns cleaned text from a .txt file
    """

    with open(txt_path, "r", encoding=encoding, errors="ignore") as file:
        text = file.read()

    return text.strip()



def extract_document_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_pdf_pages(file_path)

    elif ext in [".ppt", ".pptx"]:
        return extract_ppt_text(file_path)

    elif ext == ".txt":
        return extract_txt_text(file_path)
    
    else:
        raise ValueError(f"Unsupported file type: {ext}")


model = ChatOpenAI(
    model="gpt-4o-mini",
    api_key=os.getenv("OPENAI_API_KEY"),
    temperature=0
)

# ------------- STRUCTURED MODEL SCHEMA DEFINITION -------------------------------------------------

class VectorDbSchema(TypedDict):
    service_identity= Annotated[str,"Provide the primary service overview from the given document"]
    service_capabilities = Annotated[str,"Provide the service description including the feature description and functional capabilities"]
    service_technologies = Annotated[str,"Provide the tech stack, frameworks , cloud platforms used in this service"]
    service_intent = Annotated[str,"Provide the pbusiness outcomes and goals achieved through this service"]
    service_problems = Annotated[str,"Provide the difficulties, problems does this service solves"]
    service_industries = Annotated[str,"Provide the industry specific use cases"]
    service_usecases = Annotated[str,"Provide the case studies and usecases of this service"]
    service_figures = Annotated[str,"Provide the performance metrics,benchmarks, competitive advantages"]
    service_constraints = Annotated[str,"describe the explicit limitations, prerequisites, assumptions, and conditions under which a service is effective, applicable, or not recommended."]

structured_model = model.with_structured_output(VectorDbSchema)

# -------------------------------- PROMPT FUNCTION ---------------------------------------------

def vector_db_schema_from_doctext(text:str)->VectorDbSchema:
    prompt = PromptTemplate(
        template="""

    You are given a company service document.

    Your task is to extract structured, factual information strictly from the document text and map it to the following fields.
    Do NOT infer, assume, or add information that is not explicitly stated.
    If a field is not clearly mentioned in the document, return:
    "Not specified in the document."
    
    GENERAL RULES:
    - Be concise and factual.
    - Do not infer new keys by own.Strictly follow the schmea keys mentioned.
    - Use plain text (no markdown, no bullet symbols).
    - Avoid marketing language.
    - Do not repeat the same sentence across multiple fields.
    - Each field should contain 3-4 clear sentences unless the document provides less information.
    - Make sure for every field no information should be lef to capture.

    FIELD DEFINITIONS:

    1. service_identity  
    Provide the primary service overview from the document.  
    Explain what the service is, what it offers at a high level, and who it is for.

    2. service_capabilities  
    Provide the service description including feature descriptions and functional capabilities.  
    Focus on what the service can actually do operationally or functionally.

    3. service_technologies  
    Provide the technologies used to implement the service.  
    Include tech stack, frameworks, platforms, cloud services, tools, or AI models explicitly mentioned.

    4. service_intent  
    Provide the business outcomes and goals achieved through this service.  
    Explain why the service exists and what business value or outcomes it aims to deliver.

    5. service_problems  
    Provide the difficulties or problems that this service solves.  
    Focus on pain points, inefficiencies, or challenges explicitly addressed.

    6. service_industries  
    Provide the industry-specific applicability or use cases.  
    Mention industries, domains, or sectors where the service is applied.

    7. service_usecases  
    Provide real or described use cases, scenarios, or case studies from the document.  
    Explain how the service is used in practice.

    8. service_figures  
    Provide performance metrics, benchmarks, measurable outcomes, or competitive advantages if stated.  
    Only include quantified or clearly stated differentiators.

    9. service_constraints  
    Describe the explicit limitations, prerequisites, assumptions, or conditions under which the service is effective, applicable, or not recommended.  
    Include technical, operational, data, regulatory, or scale-related constraints only if explicitly mentioned.

    DOCUMENT TEXT:
    {DOCUMENT_TEXT}
    """,
        input_variables=["DOCUMENT_TEXT"]
    )
    response = structured_model.invoke(prompt.format(DOCUMENT_TEXT=text))
    return response

# ---------------- NORMALIZE LIST TO STR TO PERFORM STRIP FUNCTION ------------------------
def normalize_text(text) -> str:
    if isinstance(text, list):
        return "\n".join(t.strip() for t in text if t and t.strip())
    if isinstance(text, str):
        return text.strip()
    return ""



# doc_text = extract_pdf_pages('C:/Users/Navvidha_Bharech/Desktop/INT Business Central Project/BACKEND/VECTOR_DB/INT_Services_KB/INT.Data_Labs_Capability.pdf')
# output = vector_db_schema_from_doctext(doc_text)

# # Write to file
# with open("service_summary.json", "w", encoding="utf-8") as f:
#     json.dump(output, f, indent=2, ensure_ascii=False)

# print("Json save completed")

def process_documents(
    input_folder: str,
    output_folder: str = "vector_db_json_schema"
):
    os.makedirs(output_folder, exist_ok=True)

    for filename in os.listdir(input_folder):
        file_path = os.path.join(input_folder, filename)

        if not os.path.isfile(file_path):
            continue

        ext = os.path.splitext(filename)[1].lower()
        if ext not in [".pdf", ".ppt", ".pptx",".txt"]:
            continue

        print(f"Processing: {filename}")

        try:
            raw_text = extract_document_text(file_path)
            document_text = normalize_text(raw_text)

            if not document_text.strip():
                print(f"Skipping empty document: {filename}")
                continue

            structured_output = vector_db_schema_from_doctext(document_text)

            output_filename = os.path.splitext(filename)[0] + ".json"
            output_path = os.path.join(output_folder, output_filename)

            with open(output_path, "w", encoding="utf-8") as f:
                json.dump(structured_output, f, indent=2, ensure_ascii=False)

            print(f"Saved → {output_path}")

        except Exception as e:
            print(f"Failed processing {filename}: {e}")


if __name__ == "__main__":

    process_documents(input_folder="C:/Users/Navvidha_Bharech/Desktop/INT Business Central Project/BACKEND/VECTOR_DB/INT_WEBSITE_SCRAPE")
