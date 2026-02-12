import os
import json
import certifi
import pymupdf
from pptx import Presentation
from dotenv import load_dotenv
from typing import TypedDict, Annotated
from langchain_openai import ChatOpenAI
from langchain_core.prompts import PromptTemplate

load_dotenv()
os.environ["SSL_CERT_FILE"] = certifi.where()


# -------------------------------------------------------------------
# PRIMARY FUNCTION TO CONVERT PDFS INTO VECTOR DB SCHEMA ORIENTED JSON
# --------------------------------------------------------------------


# ------------- PDF LOADER ----------------------

def extract_pdf_pages(pdf_path: str) -> list[str]:
    """
    Extracts clean text from each page of a PDF.
    Returns a list of page texts.
    """
    doc = pymupdf.open(pdf_path)
    pages = []

    for page in doc:
        text = page.get_text("text")
        
        if text.strip():
            pages.append(text)

    return pages


# ------------- PPT LOADER ----------------------

def extract_ppt_text(ppt_path: str) -> str:
    prs = Presentation(ppt_path)
    slides_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slides_text.append(shape.text)

    return "\n".join(slides_text)


# ---------------- TEXT LOADER ------------------------------
def extract_txt_text(txt_path: str, encoding: str = "utf-8") -> str:
    """
    Extracts and returns cleaned text from a .txt file
    """

    with open(txt_path, "r", encoding=encoding, errors="ignore") as file:
        text = file.read()

    return text.strip()



def extract_document_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_pdf_pages(file_path)

    elif ext in [".ppt", ".pptx"]:
        return extract_ppt_text(file_path)

    elif ext == ".txt":
        return extract_txt_text(file_path)
    
    else:
        raise ValueError(f"Unsupported file type: {ext}")


model = ChatOpenAI(
    model="gpt-4o-mini",
    api_key=os.getenv("OPENAI_API_KEY"),
    temperature=0
)

# -------------- STRUCTURED MODEL SCHEMA DEFINITION -------------------------------------------------

class VectorDbSchema(TypedDict):
    service_figures= Annotated[str,"Provide the INT company's experence for delivering the following service"]
    
structured_model = model.with_structured_output(VectorDbSchema)

# -------------------------------- PROMPT FUNCTION ---------------------------------------------

def vector_db_schema_from_doctext(text:str)->VectorDbSchema:
    prompt = PromptTemplate(
        template="""

    You are given a company service document.

    Your task is to extract structured, factual information strictly from the document text and map it to the following fields.
    Do NOT infer, assume, or add information that is not explicitly stated.
    If a field is not clearly mentioned in the document, return:
    "Not specified in the document."
    
    GENERAL RULES:
    - Be concise and factual.
    - Do not infer new keys by own.Strictly follow the schmea keys mentioned.
    - Use plain text (no markdown, no bullet symbols).
    - Avoid marketing language.
    - Do not repeat the same sentence across multiple fields.
    - Each field should contain 3-4 clear sentences unless the document provides less information.
    - Make sure for every field no information should be lef to capture.

    FIELD DEFINITIONS:
    service_name :
    Create a service name after analyze the whole document

    DOCUMENT TEXT:
    {DOCUMENT_TEXT}
    """,
        input_variables=["DOCUMENT_TEXT"]
    )
    response = structured_model.invoke(prompt.format(DOCUMENT_TEXT=text))
    return response

# ---------------- NORMALIZE LIST TO STR TO PERFORM STRIP FUNCTION ------------------------
def normalize_text(text) -> str:
    if isinstance(text, list):
        return "\n".join(t.strip() for t in text if t and t.strip())
    if isinstance(text, str):
        return text.strip()
    return ""



# doc_text = extract_pdf_pages('C:/Users/Navvidha_Bharech/Desktop/INT Business Central Project/BACKEND/VECTOR_DB/INT_Services_KB/INT.Data_Labs_Capability.pdf')
# output = vector_db_schema_from_doctext(doc_text)

# # Write to file
# with open("service_summary.json", "w", encoding="utf-8") as f:
#     json.dump(output, f, indent=2, ensure_ascii=False)

# print("Json save completed")

def process_documents(
    input_folder: str,
    output_folder: str = "INT_experience_json_schema"
):
    os.makedirs(output_folder, exist_ok=True)

    for filename in os.listdir(input_folder):
        file_path = os.path.join(input_folder, filename)

        if not os.path.isfile(file_path):
            continue

        ext = os.path.splitext(filename)[1].lower()
        if ext not in [".pdf", ".ppt", ".pptx",".txt"]:
            continue

        print(f"Processing: {filename}")

        try:
            raw_text = extract_document_text(file_path)
            document_text = normalize_text(raw_text)

            if not document_text.strip():
                print(f"Skipping empty document: {filename}")
                continue

            structured_output = vector_db_schema_from_doctext(document_text)

            output_filename = os.path.splitext(filename)[0] + ".json"
            output_path = os.path.join(output_folder, output_filename)

            with open(output_path, "w", encoding="utf-8") as f:
                json.dump(structured_output, f, indent=2, ensure_ascii=False)

            print(f"Saved → {output_path}")

        except Exception as e:
            print(f"Failed processing {filename}: {e}")


if __name__ == "__main__":

    process_documents(input_folder="C:/Users/Navvidha_Bharech/Desktop/INT Business Central Project/BACKEND/VECTOR_DB/INT_Services_KB")
